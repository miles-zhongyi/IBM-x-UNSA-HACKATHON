"""
Document parser layer — converts raw files to plain text for ingest_text().

Supported formats:
  PDF    — text-based (PyMuPDF direct), scanned/handwritten (pdf2image + Tesseract OCR)
  DOCX   — python-docx direct extraction
  Images — .jpg / .jpeg / .png (Tesseract OCR, with preprocessing for handwriting)
  Excel  — .xlsx / .xls (openpyxl, row-by-row text)
  PPT    — .pptx (python-pptx, slide text extraction)

Usage:
    from parser import parse_file
    text = parse_file("/path/to/note.pdf")
    result = ingest_text(text, patient_id="p001", source_label="note.pdf")

Install dependencies:
    pip install pymupdf pdf2image pytesseract pillow python-docx openpyxl python-pptx
    # system: apt install tesseract-ocr   OR   brew install tesseract
"""

from __future__ import annotations

import os
import re
import tempfile
from pathlib import Path

# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def parse_file(filepath: str | Path, lang: str = "eng") -> str:
    """
    Parse any supported document type to plain text.

    Args:
        filepath: Path to the file on disk.
        lang:     Tesseract language code for OCR (e.g. "eng", "spa", "fra").
                  For multilingual docs, use "eng+spa" etc.

    Returns:
        Extracted text string, ready for ingest_text().

    Raises:
        ValueError: If the file type is unsupported or file does not exist.
    """
    path = Path(filepath)
    if not path.exists():
        raise ValueError(f"File not found: {filepath}")

    ext = path.suffix.lower()

    if ext == ".pdf":
        return _parse_pdf(path, lang=lang)
    elif ext == ".docx":
        return _parse_docx(path)
    elif ext in (".jpg", ".jpeg", ".png"):
        return _parse_image(path, lang=lang)
    elif ext in (".xlsx", ".xls"):
        return _parse_excel(path)
    elif ext == ".pptx":
        return _parse_pptx(path)
    else:
        raise ValueError(
            f"Unsupported file type: {ext!r}. "
            "Supported: .pdf, .docx, .jpg, .jpeg, .png, .xlsx, .xls, .pptx"
        )


def parse_bytes(data: bytes, filename: str, lang: str = "eng") -> str:
    """
    Parse a file from raw bytes (e.g. from a FastAPI UploadFile).
    Writes to a temp file, parses, then cleans up.
    """
    suffix = Path(filename).suffix.lower()
    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        tmp.write(data)
        tmp_path = tmp.name
    try:
        return parse_file(tmp_path, lang=lang)
    finally:
        os.unlink(tmp_path)


# ---------------------------------------------------------------------------
# PDF parser — auto-detects text layer vs scanned
# ---------------------------------------------------------------------------

def _parse_pdf(path: Path, lang: str = "eng") -> str:
    """
    Try direct text extraction first (fast, accurate).
    Fall back to OCR page-by-page if text layer is missing or too sparse.
    """
    try:
        import fitz  # PyMuPDF
    except ImportError:
        raise ImportError("Install PyMuPDF: pip install pymupdf")

    doc = fitz.open(str(path))
    pages_text: list[str] = []

    for page in doc:
        text = page.get_text().strip()

        # Heuristic: if direct extraction yields < 50 chars, the page is
        # likely scanned — run OCR on it instead.
        if len(text) >= 50:
            pages_text.append(text)
        else:
            ocr_text = _ocr_pdf_page(path, page.number, lang=lang)
            pages_text.append(ocr_text)

    doc.close()
    return _clean_text("\n\n".join(pages_text))


def _ocr_pdf_page(path: Path, page_number: int, lang: str) -> str:
    """Rasterize one PDF page and run Tesseract on it."""
    try:
        from pdf2image import convert_from_path
    except ImportError:
        raise ImportError("Install pdf2image: pip install pdf2image")

    images = convert_from_path(
        str(path),
        dpi=300,                         # high DPI improves OCR accuracy
        first_page=page_number + 1,      # pdf2image is 1-indexed
        last_page=page_number + 1,
    )
    if not images:
        return ""
    return _ocr_image(images[0], lang=lang, is_handwriting=False)


# ---------------------------------------------------------------------------
# Image parser — printed text and handwriting
# ---------------------------------------------------------------------------

def _parse_image(path: Path, lang: str = "eng") -> str:
    """OCR a .jpg / .png image, auto-detecting handwriting vs printed text."""
    try:
        from PIL import Image
    except ImportError:
        raise ImportError("Install Pillow: pip install pillow")

    img = Image.open(str(path))
    text = _ocr_image(img, lang=lang, is_handwriting=_looks_like_handwriting(img))
    return _clean_text(text)


def _ocr_image(img, lang: str, is_handwriting: bool) -> str:
    """Run Tesseract OCR on a PIL Image object."""
    try:
        import pytesseract
        from PIL import ImageFilter, ImageOps
    except ImportError:
        raise ImportError("Install pytesseract + Pillow: pip install pytesseract pillow")

    # Preprocess for better accuracy
    img = img.convert("L")              # grayscale
    img = ImageOps.autocontrast(img)    # normalize contrast
    img = img.filter(ImageFilter.SHARPEN)

    if is_handwriting:
        # LSTM engine (--oem 1) with no page segmentation assumption (--psm 6)
        # works better for handwritten notes
        config = "--oem 1 --psm 6"
    else:
        # Default: LSTM engine, single uniform block of text
        config = "--oem 3 --psm 6"

    return pytesseract.image_to_string(img, lang=lang, config=config)


def _looks_like_handwriting(img) -> bool:
    """
    Rough heuristic: handwritten images tend to have lower pixel std deviation
    than clean printed text. Replace with a watsonx vision call for production.
    """
    try:
        import numpy as np
        from PIL import ImageOps
        gray = ImageOps.grayscale(img)
        arr = np.array(gray, dtype=float)
        return arr.std() < 55.0
    except Exception:
        return False  # default to printed-text OCR if numpy unavailable


# ---------------------------------------------------------------------------
# DOCX parser
# ---------------------------------------------------------------------------

def _parse_docx(path: Path) -> str:
    """Extract all paragraph and table text from a .docx file."""
    try:
        import docx
    except ImportError:
        raise ImportError("Install python-docx: pip install python-docx")

    doc = docx.Document(str(path))
    parts: list[str] = []

    # Paragraphs (headings, body text)
    for para in doc.paragraphs:
        t = para.text.strip()
        if t:
            parts.append(t)

    # Tables — common for lab results in Word docs
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(
                cell.text.strip() for cell in row.cells if cell.text.strip()
            )
            if row_text:
                parts.append(row_text)

    return _clean_text("\n\n".join(parts))


# ---------------------------------------------------------------------------
# Excel parser
# ---------------------------------------------------------------------------

def _parse_excel(path: Path) -> str:
    """
    Convert Excel rows to readable text.
    Each sheet is labelled; each row is tab-joined.
    Works for lab result exports, medication lists, etc.
    """
    try:
        import openpyxl
    except ImportError:
        raise ImportError("Install openpyxl: pip install openpyxl")

    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    parts: list[str] = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f"Sheet: {sheet_name}")
        for row in ws.iter_rows(values_only=True):
            row_text = "\t".join(str(v) for v in row if v is not None)
            if row_text.strip():
                parts.append(row_text)

    wb.close()
    return _clean_text("\n".join(parts))


# ---------------------------------------------------------------------------
# PowerPoint parser
# ---------------------------------------------------------------------------

def _parse_pptx(path: Path) -> str:
    """Extract text from all slides in a .pptx file."""
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("Install python-pptx: pip install python-pptx")

    prs = Presentation(str(path))
    parts: list[str] = []

    for i, slide in enumerate(prs.slides, 1):
        slide_parts: list[str] = [f"Slide {i}:"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        slide_parts.append(t)
        if len(slide_parts) > 1:
            parts.append("\n".join(slide_parts))

    return _clean_text("\n\n".join(parts))


# ---------------------------------------------------------------------------
# Shared text cleaning
# ---------------------------------------------------------------------------

def _clean_text(text: str) -> str:
    """Normalize whitespace and strip junk characters from extracted text."""
    if not text:
        return ""
    text = re.sub(r"\n{3,}", "\n\n", text)           # collapse excess blank lines
    text = re.sub(r"[ \t]+", " ", text)               # normalize spaces/tabs
    text = re.sub(r"[^\x09\x0A\x0D\x20-\x7E\u00A0-\uFFFF]", "", text)  # strip non-printable
    return text.strip()
